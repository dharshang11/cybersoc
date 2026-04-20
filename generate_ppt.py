#!/usr/bin/env python3
"""
CyberSOC — Professional PowerPoint Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── COLOR PALETTE ───
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x12, 0x12, 0x1A)
BG_CARD2 = RGBColor(0x18, 0x18, 0x24)
RED_NEON = RGBColor(0xFF, 0x1A, 0x1A)
RED_DARK = RGBColor(0x8B, 0x00, 0x00)
RED_SOFT = RGBColor(0xFF, 0x44, 0x44)
GREEN = RGBColor(0x00, 0xFF, 0x88)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x99)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
TRANSPARENT = RGBColor(0x16, 0x16, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=RED_NEON):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        
        run1 = p.add_run()
        run1.text = "▸ "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = bullet_color
        run1.font.bold = True
        run1.font.name = "Calibri"
        
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Calibri"
    return txBox

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x00, 0x00)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = RED_NEON
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if r % 2 == 0 else BG_CARD2
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = LIGHT_GRAY
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
    
    # Remove table borders
    for r in range(len(rows) + 1):
        for c in range(len(headers)):
            cell = table.cell(r, c)
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    
    return table_shape

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED_NEON
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def slide_number(slide, num, total=15):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num}/{total}", font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Accent decorations
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)
add_shape(slide, Inches(0), Inches(7.45), Inches(13.333), Pt(4), RED_DARK)

# Left vertical accent
add_shape(slide, Inches(1.5), Inches(1.5), Pt(4), Inches(4.5), RED_NEON)

add_text(slide, Inches(2.2), Inches(1.8), Inches(9), Inches(1.2),
         "🛡️  CyberSOC", font_size=52, color=RED_NEON, bold=True)
add_text(slide, Inches(2.2), Inches(3.0), Inches(9), Inches(1),
         "AI-Powered Security Operations Center", font_size=30, color=WHITE, bold=False)

add_accent_line(slide, Inches(2.2), Inches(4.0), Inches(3))

add_text(slide, Inches(2.2), Inches(4.3), Inches(9), Inches(0.8),
         "Real-Time Threat Detection  •  3D Visualization  •  Gesture Control",
         font_size=18, color=GRAY)

add_text(slide, Inches(2.2), Inches(5.5), Inches(9), Inches(0.5),
         "A Final Year Cybersecurity Project", font_size=16, color=CYAN, bold=True)

add_text(slide, Inches(2.2), Inches(6.2), Inches(9), Inches(0.5),
         "Technologies: Python • FastAPI • Next.js • Three.js • MongoDB • MediaPipe",
         font_size=14, color=GRAY)

slide_number(slide, 1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "⚡  PROBLEM STATEMENT", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Stats cards
stats = [
    ("2,365", "Cyberattacks\nper day globally", RED_SOFT),
    ("$4.88M", "Average cost of\na data breach", ORANGE),
    ("277 days", "Average time to\nidentify a breach", YELLOW),
    ("68%", "Breaches involve\nhuman error", CYAN),
]

for i, (val, label, color) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    card = add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
    add_text(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.7),
             val, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.4), Inches(2.4), Inches(0.7),
             label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.2), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
add_text(slide, Inches(1.2), Inches(4.0), Inches(11), Inches(0.6),
         "The Challenge", font_size=22, color=WHITE, bold=True)

add_bullet_list(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(2.8), [
    "Traditional security tools are reactive — they detect AFTER damage is done",
    "SOC analysts face alert fatigue with thousands of logs per day",
    "No real-time visualization of system health and active threats",
    "Manual monitoring is slow, error-prone, and cannot scale",
    "Need: An automated, AI-driven system that detects + visualizes + responds in real-time",
], font_size=16)

slide_number(slide, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: PROPOSED SOLUTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "💡  PROPOSED SOLUTION", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

features = [
    ("🕵️  REAL-TIME MONITORING", "Agent collects CPU, RAM, disk, processes,\nnetwork connections & WiFi data every 2 seconds", RED_SOFT),
    ("🧠  AI THREAT DETECTION", "5 intelligent rules detect suspicious processes,\nreverse shells, port scans & resource abuse", CYAN),
    ("🌍  3D VISUALIZATION", "Interactive Three.js globe with attack arc\nanimations showing threats in real-time", GREEN),
    ("✋  GESTURE CONTROL", "MediaPipe hand tracking lets you control\nthe dashboard with hand gestures via webcam", YELLOW),
    ("🔍  DEEP SCAN ENGINE", "On-demand full system scan: file integrity,\nprocess audit, network analysis & anomalies", ORANGE),
    ("📡  LIVE WEBSOCKET", "Instant data broadcast — zero polling delay\nfrom agent to dashboard via WebSocket", RGBColor(0xBB, 0x86, 0xFC)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.9)
    
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.6),
             title, font_size=16, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.3), Inches(1.3),
             desc, font_size=14, color=LIGHT_GRAY)

slide_number(slide, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🏗️  SYSTEM ARCHITECTURE", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Architecture diagram using shapes
# Agent box
add_shape(slide, Inches(0.5), Inches(2.0), Inches(3.2), Inches(4.0), BG_CARD, border_color=CYAN, border_width=Pt(2))
add_text(slide, Inches(0.7), Inches(2.1), Inches(2.8), Inches(0.4),
         "🖥️  AGENT (Python)", font_size=15, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(0.7), Inches(2.6), Inches(2.8), Inches(3.0), [
    "psutil → CPU/RAM/Disk",
    "psutil → Process list",
    "lsof → Network connections",
    "system_profiler → WiFi",
    "hashlib → File integrity",
    "Sends every 2 seconds",
], font_size=12, color=LIGHT_GRAY, bullet_color=CYAN)

# Arrow
add_text(slide, Inches(3.7), Inches(3.5), Inches(1.2), Inches(0.5),
         "HTTP POST ►", font_size=12, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# Backend box
add_shape(slide, Inches(4.9), Inches(1.8), Inches(3.6), Inches(4.5), BG_CARD, border_color=RED_NEON, border_width=Pt(2))
add_text(slide, Inches(5.1), Inches(1.9), Inches(3.2), Inches(0.4),
         "⚙️  BACKEND (FastAPI)", font_size=15, color=RED_NEON, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(5.1), Inches(2.4), Inches(3.2), Inches(3.5), [
    "main.py → API Server",
    "engine.py → Threat Detection",
    "models.py → Data Validation",
    "sockets.py → WS Manager",
    "database.py → MongoDB",
    "Detects threats instantly",
    "Broadcasts via WebSocket",
], font_size=12, color=LIGHT_GRAY, bullet_color=RED_NEON)

# Arrow
add_text(slide, Inches(8.5), Inches(3.5), Inches(1.2), Inches(0.5),
         "WebSocket ►", font_size=12, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Frontend box
add_shape(slide, Inches(9.7), Inches(2.0), Inches(3.2), Inches(4.0), BG_CARD, border_color=GREEN, border_width=Pt(2))
add_text(slide, Inches(9.9), Inches(2.1), Inches(2.8), Inches(0.4),
         "🌐  FRONTEND (Next.js)", font_size=15, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(9.9), Inches(2.6), Inches(2.8), Inches(3.0), [
    "page.tsx → Dashboard UI",
    "ThreatGlobe → 3D Globe",
    "GestureControl → Hands",
    "Real-time data display",
    "4 scrollable pages",
    "Dark cyberpunk theme",
], font_size=12, color=LIGHT_GRAY, bullet_color=GREEN)

# DB box at bottom
add_shape(slide, Inches(5.4), Inches(6.5), Inches(2.6), Inches(0.8), BG_CARD, border_color=ORANGE, border_width=Pt(2))
add_text(slide, Inches(5.4), Inches(6.55), Inches(2.6), Inches(0.7),
         "🗄️ MongoDB", font_size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(slide, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "📦  TECHNOLOGY STACK", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
    ["Technology", "Type", "Purpose", "Why Chosen"],
    [
        ["FastAPI (Python)", "Backend Framework", "API Server + WebSocket Hub", "Fastest Python framework; async; built-in WS support"],
        ["Next.js 16 (React 19)", "Frontend Framework", "Dashboard UI + Routing", "Server-side rendering; optimized builds; TypeScript"],
        ["Three.js", "3D Graphics", "Interactive Threat Globe", "WebGL rendering; smooth 3D animations in browser"],
        ["MongoDB + Motor", "Database", "Store alerts, logs, scans", "NoSQL flexibility; async driver; schema-free"],
        ["MediaPipe Hands", "ML / AI", "Hand Gesture Recognition", "Google's ML model; runs in-browser; real-time"],
        ["Socket.io / WebSocket", "Real-time Protocol", "Live data streaming", "Full-duplex; instant push; no polling needed"],
        ["psutil", "System Monitor", "CPU, RAM, Disk, Processes", "Cross-platform; detailed system metrics"],
        ["Pydantic", "Data Validation", "Validate agent payloads", "Type-safe schemas; auto-validation"],
        ["Serveo.net", "Tunneling", "Expose localhost to internet", "Free SSH tunnel; instant public URL"],
        ["TypeScript", "Language", "Type-safe frontend code", "Catch bugs at compile time; better DX"],
    ],
    col_widths=[Inches(2.5), Inches(2), Inches(3.2), Inches(4.6)]
)

slide_number(slide, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: DATA FLOW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🔄  DATA FLOW — How It All Connects", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

steps = [
    ("STEP 1", "COLLECT", "Agent reads CPU, RAM, disk,\nprocesses, network, WiFi\nusing psutil + lsof", CYAN, "Every 2 seconds"),
    ("STEP 2", "TRANSMIT", "Agent sends JSON payload\nvia HTTP POST to\n/api/agent/ingest", YELLOW, "With API key auth"),
    ("STEP 3", "DETECT", "Backend runs 5 threat\ndetection rules on\nthe incoming data", RED_SOFT, "engine.py rules"),
    ("STEP 4", "STORE", "Alerts + metrics saved\nto MongoDB for\nhistorical analysis", ORANGE, "Async via Motor"),
    ("STEP 5", "BROADCAST", "WebSocket pushes\nLIVE_METRICS + NEW_ALERTS\nto all dashboards", GREEN, "Instant, zero delay"),
    ("STEP 6", "VISUALIZE", "Dashboard renders 3D globe,\nlive charts, alert cards,\nand system monitors", RGBColor(0xBB, 0x86, 0xFC), "React re-renders"),
]

for i, (step, title, desc, color, note) in enumerate(steps):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(1.5 + (i // 3) * 3.0)
    
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.6), BG_CARD, border_color=color, border_width=Pt(2))
    
    add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(1.2), Inches(0.4),
             step, font_size=11, color=GRAY, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.5),
             title, font_size=22, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.1), Inches(3.5), Inches(1.0),
             desc, font_size=14, color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), y + Inches(2.1), Inches(3.5), Inches(0.4),
             f"📌 {note}", font_size=11, color=GRAY)

slide_number(slide, 6)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: THREAT DETECTION ENGINE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🧠  THREAT DETECTION ENGINE", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
         "5 intelligent detection rules analyze every agent payload in real-time:",
         font_size=16, color=GRAY)

add_table(slide, Inches(0.5), Inches(1.9), Inches(12.3), Inches(3.0),
    ["#", "Rule", "Trigger Condition", "Severity", "Example Scenario"],
    [
        ["1", "🔥 High CPU", "CPU usage > 90%", "HIGH", "Cryptominer consuming 99% CPU"],
        ["2", "💾 High RAM", "RAM usage > 90%", "MEDIUM", "Memory leak consuming all RAM"],
        ["3", "🕵️ Suspicious Process", "Process name in blacklist", "CRITICAL", "nmap, metasploit, hashcat detected"],
        ["4", "🌐 High Connections", "50+ network connections", "MEDIUM", "DDoS or network flooding attack"],
        ["5", "🚪 Suspicious Ports", "Connection to 4444, 1337, etc.", "CRITICAL", "Reverse shell / C2 communication"],
    ],
    col_widths=[Inches(0.5), Inches(2.2), Inches(3), Inches(1.5), Inches(5.1)]
)

add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
add_text(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.5),
         "🔴 Blacklisted Process Names (19 Tools)", font_size=16, color=RED_SOFT, bold=True)
add_text(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.2),
         "nmap  •  netcat  •  hashcat  •  mimikatz  •  wireshark  •  hydra  •  sqlmap  •  metasploit\nmeterpreter  •  ettercap  •  aircrack-ng  •  burpsuite  •  nikto  •  gobuster  •  ncat\nmasscan  •  responder  •  crackmapexec  •  evil-winrm",
         font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(slide, 7)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: DASHBOARD — 4 PAGES
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "📱  DASHBOARD — 4 Interactive Pages", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

pages = [
    ("PAGE 1", "🌍 THREAT MAP", [
        "3D rotating wireframe globe (Three.js)",
        "Continent dots with city coordinates",
        "Animated attack arc curves",
        "Live CPU / RAM / Disk stat bars",
        "WiFi network info panel",
        "Threat counter + connection count",
    ], GREEN),
    ("PAGE 2", "🖥️ SYSTEM MONITOR", [
        "Active process table (top 15 by CPU)",
        "WiFi details: SSID, BSSID, channel",
        "Network connections: local ↔ remote",
        "Connection status colors",
        "Real-time table updates via WS",
    ], CYAN),
    ("PAGE 3", "🔍 DEEP SCAN", [
        "On-demand system scan button",
        "Scan summary: threat level + findings",
        "Process inspection (suspicious PIDs)",
        "File integrity: SHA-256 hash checks",
        "System anomaly detection",
    ], ORANGE),
    ("PAGE 4", "🚨 ALERTS CENTER", [
        "Severity-coded alert cards",
        "Critical/High/Medium/Low colors",
        "Auto-scroll on critical threat",
        "Full-screen alert on critical detection",
        "Real-time alert stream",
    ], RED_SOFT),
]

for i, (page, title, items, color) in enumerate(pages):
    x = Inches(0.5 + i * 3.15)
    card = add_shape(slide, x, Inches(1.5), Inches(2.95), Inches(5.6), BG_CARD, border_color=color, border_width=Pt(2))
    
    add_text(slide, x + Inches(0.1), Inches(1.6), Inches(2.75), Inches(0.3),
             page, font_size=11, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(1.9), Inches(2.75), Inches(0.5),
             title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_bullet_list(slide, x + Inches(0.15), Inches(2.6), Inches(2.6), Inches(4.0),
                    items, font_size=12, color=LIGHT_GRAY, bullet_color=color)

slide_number(slide, 8)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: 3D THREAT GLOBE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🌍  3D THREAT GLOBE — Three.js Visualization", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Globe features
globe_features = [
    ("Wireframe Sphere", "40×40 segment sphere with red wireframe\nand semi-transparent material (opacity: 0.15)", "SphereGeometry(1, 40, 40)"),
    ("Inner Glow", "Dark red inner sphere creates depth effect\nwith opacity: 0.4 for volumetric feel", "SphereGeometry(0.98, 32, 32)"),
    ("Continent Dots", "55 real-world city coordinates\nmapped using lat/lon → 3D conversion", "phi = (90 - lat) × π/180"),
    ("Attack Arcs", "Quadratic Bezier curves between random\ncities with auto-expire after 3-7 seconds", "QuadraticBezierCurve3()"),
    ("Orbital Rings", "Two tilted torus rings at different angles\nrotating in opposite directions", "TorusGeometry(1.4, 0.005)"),
    ("Star Particles", "1500 random points in 3D space\nfor immersive space background", "Points + PointsMaterial"),
]

for i, (title, desc, code) in enumerate(globe_features):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.5 + row * 2.0)
    
    card = add_shape(slide, x, y, Inches(6.1), Inches(1.8), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.5), Inches(0.4),
             f"🔴 {title}", font_size=16, color=RED_SOFT, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.5), Inches(0.6),
             desc, font_size=13, color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(5.5), Inches(0.4),
             f"Code: {code}", font_size=11, color=CYAN)

slide_number(slide, 9)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: GESTURE CONTROL
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "✋  GESTURE CONTROL — MediaPipe Hand Tracking", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

gestures = [
    ("🖐️", "OPEN PALM", "Trigger Deep Scan", "All 4 fingertips above knuckles", GREEN),
    ("✊", "FIST", "Dismiss Alert / Go Home", "All 4 fingertips below knuckles", RED_SOFT),
    ("👍", "THUMBS UP", "Show Motivational Quote", "Thumb above index + fist closed", YELLOW),
    ("☝️", "INDEX FINGER", "Navigate to Alerts Page", "Index up, others curled down", CYAN),
]

for i, (emoji, name, action, detection, color) in enumerate(gestures):
    x = Inches(0.5 + i * 3.15)
    card = add_shape(slide, x, Inches(1.5), Inches(2.95), Inches(3.2), BG_CARD, border_color=color, border_width=Pt(2))
    
    add_text(slide, x, Inches(1.6), Inches(2.95), Inches(0.7),
             emoji, font_size=40, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.3), Inches(2.55), Inches(0.4),
             name, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.8), Inches(2.55), Inches(0.4),
             f"Action: {action}", font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.3), Inches(2.55), Inches(0.7),
             f"Detection: {detection}", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# How it works
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
add_text(slide, Inches(0.9), Inches(5.1), Inches(11), Inches(0.5),
         "How MediaPipe Hand Tracking Works", font_size=18, color=WHITE, bold=True)

add_bullet_list(slide, Inches(0.9), Inches(5.6), Inches(11), Inches(1.5), [
    "Webcam stream → MediaPipe Hands ML model (runs entirely in-browser, no server needed)",
    "Model outputs 21 hand landmark coordinates (x, y, z) per frame at 30fps",
    "Compare fingertip Y-position vs knuckle Y-position: tip above knuckle = finger UP",
    "Gesture matched → action triggered with 2.5-second cooldown to prevent spam",
], font_size=14)

slide_number(slide, 10)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: DEEP SCAN ENGINE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🔍  DEEP SCAN ENGINE", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

scans = [
    ("🕵️ Process Inspection", "Scans ALL running processes for:\n• Names matching 19 suspicious tools\n• CPU usage > 80% (resource hogs)\n• Memory usage > 50%\n• Reports PID, executable path, user", RED_SOFT),
    ("🌐 Network Analysis", "Analyzes all network connections:\n• Counts external connections\n• Flags suspicious ports (4444, 1337, etc.)\n• Identifies C2 communication\n• Reports process + remote address", CYAN),
    ("📂 File Integrity", "SHA-256 hash checks on critical files:\n• /etc/hosts — DNS hijacking\n• /etc/passwd — User tampering\n• ~/.ssh/authorized_keys — Backdoors\n• ~/.zshrc — Shell injection", GREEN),
    ("⚠️ Anomaly Detection", "System-wide anomaly checks:\n• Process count > 400 = suspicious\n• Swap usage > 80% = critical\n• System uptime analysis\n• Disk I/O statistics", ORANGE),
]

for i, (title, desc, color) in enumerate(scans):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.5 + row * 3.0)
    
    card = add_shape(slide, x, y, Inches(6.1), Inches(2.7), BG_CARD, border_color=color, border_width=Pt(2))
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.5), Inches(0.5),
             title, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.5), Inches(1.7),
             desc, font_size=13, color=LIGHT_GRAY)

slide_number(slide, 11)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: ATTACK SIMULATOR
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🧪  ATTACK SIMULATOR — 12 Scenarios", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_table(slide, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8),
    ["#", "Attack Name", "What It Simulates", "Alerts Triggered"],
    [
        ["1", "NMAP Scan", "Port scanning tool running", "suspicious_process (Critical)"],
        ["2", "Metasploit", "Exploitation framework active", "suspicious_process ×2 (Critical)"],
        ["3", "Reverse Shell", "Backdoor on port 4444", "suspicious_process + suspicious_port"],
        ["4", "Hashcat", "Password cracking (98% CPU)", "suspicious_process + high_cpu"],
        ["5", "Nikto", "Web vulnerability scanner", "suspicious_process (Critical)"],
        ["6", "Wireshark", "Packet sniffing/capture", "suspicious_process (Critical)"],
        ["7", "Hydra", "SSH brute force attack", "suspicious_process (Critical)"],
        ["8", "CPU Spike", "Cryptominer at 99% CPU", "high_cpu + suspicious_port"],
        ["9", "Gobuster", "Directory brute force", "suspicious_process (Critical)"],
        ["10", "SQLMap", "SQL injection attack", "suspicious_process (Critical)"],
        ["11", "C2 Ports", "Connections to 4444, 5555, 1337", "suspicious_connection ×3"],
        ["12", "DDoS Flood", "55 simultaneous connections", "high_connections (Medium)"],
    ],
    col_widths=[Inches(0.6), Inches(2.0), Inches(4.5), Inches(5.6)]
)

slide_number(slide, 12)


# ══════════════════════════════════════════════════════════════
# SLIDE 13: API ENDPOINTS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🔐  API ENDPOINTS & COMMUNICATION", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_table(slide, Inches(0.3), Inches(1.4), Inches(12.7), Inches(3.5),
    ["Method", "Endpoint", "Called By", "Purpose"],
    [
        ["POST", "/api/agent/ingest", "Agent / Simulator", "Receives system metrics, runs detection, broadcasts"],
        ["POST", "/api/agent/scan-report", "Agent", "Receives deep scan results from agent"],
        ["POST", "/api/scan/trigger", "Dashboard", "Frontend requests a deep scan"],
        ["GET", "/api/scan/status", "Dashboard", "Check if scan is currently running"],
        ["GET", "/api/scan/last-report", "Dashboard", "Fetch most recent scan report"],
        ["GET", "/api/agent/health", "Simulator", "Health check — is backend online?"],
        ["GET", "/api/alerts", "Dashboard", "Fetch alert history from MongoDB"],
        ["WS", "/ws", "Dashboard", "Real-time bidirectional WebSocket"],
    ],
    col_widths=[Inches(1.2), Inches(3.0), Inches(2.5), Inches(6.0)]
)

# Communication protocols
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x3A), border_width=Pt(1))
add_text(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4),
         "Communication Protocols", font_size=18, color=WHITE, bold=True)

protocols = [
    ("Agent → Backend", "HTTP POST", "JSON payload every 2s", CYAN),
    ("Backend → MongoDB", "MongoDB Wire Protocol", "BSON documents via Motor async", ORANGE),
    ("Backend ↔ Dashboard", "WebSocket", "Full-duplex JSON messages, instant push", GREEN),
    ("Dashboard (Camera)", "WebRTC", "Video frames processed by MediaPipe ML", YELLOW),
]

for i, (route, protocol, desc, color) in enumerate(protocols):
    x = Inches(0.9 + i * 3.1)
    add_text(slide, x, Inches(5.8), Inches(2.8), Inches(0.3),
             route, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.1), Inches(2.8), Inches(0.3),
             protocol, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.5), Inches(2.8), Inches(0.4),
             desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

slide_number(slide, 13)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: HOW TO RUN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "🚀  HOW TO RUN THE PROJECT", font_size=28, color=RED_NEON, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

terminals = [
    ("Terminal 1", "🔧 Start Backend", "cd backend\npip install -r requirements.txt\nuvicorn main:app --reload --port 8000", RED_SOFT),
    ("Terminal 2", "🌐 Start Frontend", "cd frontend\nnpm install\nnpm run dev", GREEN),
    ("Terminal 3", "🕵️ Start Agent", "cd agent\npip install psutil requests\npython3 agent.py", CYAN),
    ("Terminal 4", "💣 Attack Simulator", "python3 attack_simulator.py\n\n(Injects 12 fake attacks\nto test alert detection)", ORANGE),
]

for i, (term, title, cmds, color) in enumerate(terminals):
    x = Inches(0.5 + i * 3.15)
    card = add_shape(slide, x, Inches(1.4), Inches(2.95), Inches(3.3), BG_CARD, border_color=color, border_width=Pt(2))
    
    add_text(slide, x + Inches(0.1), Inches(1.5), Inches(2.75), Inches(0.3),
             term, font_size=11, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(1.8), Inches(2.75), Inches(0.4),
             title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), Inches(2.3), Inches(2.65), Inches(2.2),
             cmds, font_size=12, color=LIGHT_GRAY)

# Share online
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), BG_CARD, border_color=RGBColor(0xBB, 0x86, 0xFC), border_width=Pt(2))
add_text(slide, Inches(0.9), Inches(5.1), Inches(11), Inches(0.5),
         "🌐  Share With Friends (Serveo.net Tunnel)", font_size=18, color=RGBColor(0xBB, 0x86, 0xFC), bold=True)
add_text(slide, Inches(0.9), Inches(5.7), Inches(5.5), Inches(1.2),
         "# Tunnel Backend (Terminal 5)\nssh -R 80:localhost:8000 serveo.net\n→ https://xxxxx.serveo.net\n\n# Tunnel Frontend (Terminal 6)\nssh -R 80:localhost:3000 serveo.net\n→ https://yyyyy.serveo.net  ← Share this!",
         font_size=13, color=LIGHT_GRAY)
add_text(slide, Inches(7), Inches(5.7), Inches(5.5), Inches(1.2),
         "Your friend just opens the link in a browser!\n\nTo monitor their own system:\n1. Install Python + psutil\n2. Run agent.py pointing to your backend URL\n3. Their system data appears on your dashboard",
         font_size=13, color=LIGHT_GRAY)

slide_number(slide, 14)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: THANK YOU
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), RED_NEON)
add_shape(slide, Inches(0), Inches(7.45), Inches(13.333), Pt(4), RED_DARK)

add_shape(slide, Inches(1.5), Inches(1.5), Pt(4), Inches(4.5), RED_NEON)

add_text(slide, Inches(2.2), Inches(2.0), Inches(9), Inches(1),
         "Thank You!", font_size=52, color=RED_NEON, bold=True)

add_accent_line(slide, Inches(2.2), Inches(3.2), Inches(3))

add_text(slide, Inches(2.2), Inches(3.6), Inches(9), Inches(0.6),
         "CyberSOC — AI-Powered Security Operations Center", font_size=22, color=WHITE)

add_text(slide, Inches(2.2), Inches(4.5), Inches(9), Inches(2.0),
         "Key Highlights:\n• Real-time system monitoring via Python agent\n• 5-rule AI threat detection engine\n• 3D Three.js interactive threat globe\n• MediaPipe gesture control (4 gestures)\n• 12 attack simulation scenarios\n• WebSocket live data streaming\n• MongoDB persistent storage",
         font_size=16, color=LIGHT_GRAY)

add_text(slide, Inches(2.2), Inches(6.5), Inches(9), Inches(0.5),
         "\"Security is not a product, but a process.\" — Bruce Schneier",
         font_size=14, color=GRAY)

slide_number(slide, 15)


# ─── SAVE ───
output_path = os.path.expanduser("~/Desktop/cyber/CyberSOC_Presentation.pptx")
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"   Total slides: 15")
print(f"   Open with PowerPoint, Keynote, or Google Slides")
